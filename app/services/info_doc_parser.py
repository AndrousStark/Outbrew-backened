"""
Intelligent Info Doc Parser Service - ULTRA POWERFUL VERSION

Ultra-intelligent company/service info document parser with:
- PDF, DOCX, and PPTX support
- 400+ heading variations across multiple languages
- Intelligent product/service extraction
- Pricing detection with multiple currency support
- Benefits and value proposition extraction
- Contact information parsing
- Team member detection
- Advanced section extraction with confidence scoring
- Flexible format handling
"""

import re
import logging
from typing import Dict, List, Optional, Tuple, Any, Set
from dataclasses import dataclass, field
try:
    import pdfplumber
except ImportError:
    pdfplumber = None

try:
    from docx import Document
except ImportError:
    Document = None

try:
    from fuzzywuzzy import fuzz
except ImportError:
    from rapidfuzz import fuzz
import io

logger = logging.getLogger(__name__)


# ==================== DATA CLASSES ====================

@dataclass
class ProductService:
    """Extracted product or service"""
    name: str
    description: Optional[str] = None
    pricing: Optional[str] = None
    features: List[str] = field(default_factory=list)
    category: Optional[str] = None  # product, service, solution, platform


@dataclass
class PricingTier:
    """Extracted pricing tier"""
    name: str
    price: str
    billing_cycle: Optional[str] = None  # monthly, yearly, one-time
    currency: Optional[str] = None
    features: List[str] = field(default_factory=list)
    is_popular: bool = False
    is_enterprise: bool = False


@dataclass
class TeamMember:
    """Extracted team member"""
    name: str
    role: Optional[str] = None
    title: Optional[str] = None
    email: Optional[str] = None
    linkedin: Optional[str] = None


@dataclass
class ContactInfo:
    """Extracted contact information"""
    emails: List[str] = field(default_factory=list)
    phones: List[str] = field(default_factory=list)
    websites: List[str] = field(default_factory=list)
    social_media: Dict[str, str] = field(default_factory=dict)  # platform: url
    address: Optional[str] = None


@dataclass
class InfoDocSection:
    """Parsed info doc section with confidence"""
    section_type: str
    heading: str
    content: str
    confidence: float  # 0-100
    line_start: int = 0
    line_end: int = 0


@dataclass
class ParsedInfoDoc:
    """Complete parsed info doc data"""
    # Company Info
    company_name: Optional[str] = None
    tagline: Optional[str] = None
    industry: Optional[str] = None

    # Products & Services
    products_services: List[ProductService] = field(default_factory=list)

    # Value Proposition
    key_benefits: List[str] = field(default_factory=list)
    unique_selling_points: List[str] = field(default_factory=list)  # Includes differentiators
    problem_solved: Optional[str] = None

    # Pricing
    pricing_tiers: List[PricingTier] = field(default_factory=list)
    pricing_raw: Optional[str] = None

    # Contact
    contact_info: ContactInfo = field(default_factory=ContactInfo)

    # Team
    team_members: List[TeamMember] = field(default_factory=list)

    # Metadata
    raw_text: str = ""
    detected_sections: List[InfoDocSection] = field(default_factory=list)
    confidence_score: float = 0.0
    warnings: List[str] = field(default_factory=list)
    word_count: int = 0

    def to_dict(self) -> Dict[str, Any]:
        """Convert to dictionary for JSON serialization"""
        return {
            "company_name": self.company_name,
            "tagline": self.tagline,
            "industry": self.industry,
            "products_services": [
                {
                    "name": p.name,
                    "description": p.description,
                    "pricing": p.pricing,
                    "features": p.features,
                    "category": p.category
                }
                for p in self.products_services
            ],
            "key_benefits": self.key_benefits,
            "unique_selling_points": self.unique_selling_points,
            "problem_solved": self.problem_solved,
            "pricing_tiers": [
                {
                    "name": t.name,
                    "price": t.price,
                    "billing_cycle": t.billing_cycle,
                    "currency": t.currency,
                    "features": t.features,
                    "is_popular": t.is_popular,
                    "is_enterprise": t.is_enterprise
                }
                for t in self.pricing_tiers
            ],
            "pricing_raw": self.pricing_raw,
            "contact_info": {
                "emails": self.contact_info.emails,
                "phones": self.contact_info.phones,
                "websites": self.contact_info.websites,
                "social_media": self.contact_info.social_media,
                "address": self.contact_info.address
            },
            "team_members": [
                {
                    "name": m.name,
                    "role": m.role,
                    "title": m.title,
                    "email": m.email,
                    "linkedin": m.linkedin
                }
                for m in self.team_members
            ],
            "confidence_score": self.confidence_score,
            "warnings": self.warnings,
            "word_count": self.word_count,
            "detected_sections": [
                {
                    "type": s.section_type,
                    "heading": s.heading,
                    "confidence": s.confidence,
                    "preview": s.content[:150] + "..." if len(s.content) > 150 else s.content
                }
                for s in self.detected_sections
            ]
        }


# ==================== MAIN PARSER CLASS ====================

class IntelligentInfoDocParser:
    """
    Ultra-intelligent info doc parser - POWERFUL VERSION

    Features:
    - 400+ heading variations
    - Multi-currency pricing detection
    - Product/service extraction
    - Benefits and USP parsing
    - Team member detection
    - Multi-language support
    """

    # ==================== SECTION PATTERNS (400+ variations) ====================

    # Company/About patterns - 60+ variations
    COMPANY_PATTERNS = [
        # English
        "about us", "about the company", "company overview", "who we are", "our story",
        "company profile", "our company", "company information", "about", "introduction",
        "company history", "background", "our background", "company background",
        "our mission", "mission", "our vision", "vision", "mission and vision",
        "mission statement", "vision statement", "about our company", "who are we",
        "company description", "what we do", "overview", "executive summary",
        "company summary", "business overview", "organization overview", "firm overview",
        "about the firm", "about the organization", "about this company", "company intro",
        "brief introduction", "company brief", "corporate overview", "corporate profile",
        "business profile", "organization profile", "firm profile", "enterprise overview",
        # German
        "über uns", "unternehmen", "firmenprofil", "unternehmensübersicht", "wer wir sind",
        # French
        "à propos", "qui sommes-nous", "notre entreprise", "présentation", "notre histoire",
        # Spanish
        "sobre nosotros", "quiénes somos", "nuestra empresa", "nuestra historia",
        # Italian
        "chi siamo", "la nostra azienda", "profilo aziendale",
        # Portuguese
        "sobre nós", "quem somos", "nossa empresa",
        # Dutch
        "over ons", "wie zijn wij", "bedrijfsprofiel",
        # Common typos
        "abot us", "abut us", "about uss", "compnay overview"
    ]

    # Products/Services patterns - 50+ variations
    PRODUCTS_PATTERNS = [
        # English
        "products", "services", "our products", "our services", "offerings", "our offerings",
        "what we offer", "solutions", "our solutions", "product line", "service offerings",
        "capabilities", "what we do", "features", "product features", "service features",
        "products and services", "services and products", "product catalog", "service catalog",
        "product portfolio", "service portfolio", "our portfolio", "product range",
        "service range", "product suite", "service suite", "platform", "our platform",
        "tools", "our tools", "applications", "our applications", "software", "our software",
        "products we offer", "services we provide", "what we provide", "deliverables",
        "product offerings", "solutions we offer", "core services", "core products",
        "main services", "main products", "key services", "key products", "flagship products",
        # German
        "produkte", "dienstleistungen", "unsere produkte", "unsere dienstleistungen", "lösungen",
        # French
        "produits", "services", "nos produits", "nos services", "solutions",
        # Spanish
        "productos", "servicios", "nuestros productos", "nuestros servicios", "soluciones",
        # Italian
        "prodotti", "servizi", "i nostri prodotti", "i nostri servizi",
        # Common typos
        "produts", "servies", "producst", "serivces"
    ]

    # Pricing patterns - 40+ variations
    PRICING_PATTERNS = [
        # English
        "pricing", "prices", "plans", "pricing plans", "packages", "our plans",
        "subscription", "subscriptions", "cost", "costs", "rates", "our rates",
        "fees", "pricing tiers", "plans and pricing", "pricing and plans",
        "investment", "pricing options", "plan options", "membership", "memberships",
        "pricing structure", "price list", "rate card", "pricing table", "plan comparison",
        "choose your plan", "select a plan", "pick a plan", "get started", "start now",
        "pricing details", "cost breakdown", "pricing breakdown", "how much", "what it costs",
        "affordable pricing", "competitive pricing", "transparent pricing", "simple pricing",
        "flexible pricing", "custom pricing", "enterprise pricing", "business pricing",
        # German
        "preise", "preisliste", "pakete", "tarife", "abonnements",
        # French
        "tarifs", "prix", "forfaits", "abonnements",
        # Spanish
        "precios", "tarifas", "planes", "paquetes",
        # Italian
        "prezzi", "listino prezzi", "piani", "abbonamenti",
        # Common typos
        "pricng", "pircing", "priceing", "plnas"
    ]

    # Benefits patterns - 45+ variations
    BENEFITS_PATTERNS = [
        # English
        "benefits", "key benefits", "advantages", "why choose us", "why us",
        "value proposition", "what you get", "features and benefits", "our advantages",
        "what's included", "our advantage", "unique benefits", "core benefits",
        "main benefits", "top benefits", "primary benefits", "customer benefits",
        "business benefits", "your benefits", "the benefits", "benefit highlights",
        "why choose", "reasons to choose", "why work with us", "why partner with us",
        "value we provide", "value we deliver", "what we deliver", "our value",
        "added value", "key advantages", "competitive advantages", "our strengths",
        "strengths", "highlights", "key highlights", "feature highlights",
        "what makes us different", "what sets us apart", "our difference",
        # German
        "vorteile", "ihre vorteile", "warum wir", "unsere stärken",
        # French
        "avantages", "vos avantages", "pourquoi nous choisir", "nos atouts",
        # Spanish
        "beneficios", "ventajas", "por qué elegirnos", "nuestras fortalezas",
        # Italian
        "vantaggi", "i vostri vantaggi", "perché sceglierci",
        # Common typos
        "benifits", "benfits", "advantges", "beneftis"
    ]

    # USP/Differentiators patterns - 40+ variations (merged)
    USP_PATTERNS = [
        # English
        "unique selling points", "usp", "usps", "differentiators", "what makes us unique",
        "our difference", "why we're different", "how we're different", "stand out",
        "what sets us apart", "competitive edge", "our edge", "unique features",
        "unique advantages", "distinctive features", "special features", "key differentiators",
        "our uniqueness", "unique value", "unique approach", "our approach",
        "innovation", "our innovation", "innovative features", "cutting edge",
        "industry leading", "market leading", "best in class", "world class",
        "one of a kind", "unlike others", "unlike competitors", "the difference",
        "differentiating factors", "competitive differentiators", "unique capabilities",
        "proprietary", "proprietary technology", "patented", "exclusive features",
        # German
        "alleinstellungsmerkmale", "was uns unterscheidet", "unsere besonderheiten",
        # French
        "ce qui nous distingue", "nos points forts", "notre différence",
        # Spanish
        "lo que nos diferencia", "nuestros puntos fuertes", "nuestra diferencia",
        # Italian
        "cosa ci distingue", "i nostri punti di forza",
        # Common typos
        "diferentiators", "diferenctiators", "uniqe", "uniuqe"
    ]

    # Problem/Solution patterns - 35+ variations
    PROBLEM_PATTERNS = [
        # English
        "problem", "the problem", "problem we solve", "problems we solve",
        "challenge", "challenges", "the challenge", "pain points", "pain point",
        "issues", "the issue", "struggles", "frustrations", "difficulties",
        "solution", "the solution", "our solution", "how we help", "how we solve",
        "we solve", "solving", "addressing", "tackling", "overcoming",
        "problem statement", "challenge statement", "the opportunity", "opportunity",
        "market problem", "industry problem", "customer problem", "business problem",
        "common problems", "typical challenges", "key challenges",
        # German
        "problem", "herausforderung", "lösung", "unsere lösung",
        # French
        "problème", "défi", "solution", "notre solution",
        # Spanish
        "problema", "desafío", "solución", "nuestra solución",
        # Italian
        "problema", "sfida", "soluzione", "la nostra soluzione"
    ]

    # Team patterns - 35+ variations
    TEAM_PATTERNS = [
        # English
        "team", "our team", "the team", "leadership", "leadership team",
        "management", "management team", "founders", "co-founders", "founding team",
        "meet the team", "about the team", "key people", "our people",
        "executives", "executive team", "board", "board of directors", "advisors",
        "advisory board", "team members", "staff", "our staff", "employees",
        "the people behind", "who's behind", "our experts", "experts",
        "core team", "founding members", "key team members", "senior team",
        "leadership and team", "people", "our leadership", "meet our team",
        # German
        "team", "unser team", "führungsteam", "gründer", "mitarbeiter",
        # French
        "équipe", "notre équipe", "direction", "fondateurs", "collaborateurs",
        # Spanish
        "equipo", "nuestro equipo", "liderazgo", "fundadores",
        # Italian
        "team", "il nostro team", "leadership", "fondatori"
    ]

    # Contact patterns - 30+ variations
    CONTACT_PATTERNS = [
        # English
        "contact", "contact us", "get in touch", "reach us", "reach out",
        "connect", "connect with us", "contact information", "how to reach us",
        "talk to us", "speak to us", "let's talk", "let's connect", "get started",
        "request demo", "request a demo", "book a demo", "schedule a call",
        "schedule a meeting", "book a call", "free consultation", "contact details",
        "our contact", "contact info", "reach out to us", "drop us a line",
        "send us a message", "email us", "call us", "find us",
        # German
        "kontakt", "kontaktieren sie uns", "erreichen sie uns",
        # French
        "contact", "contactez-nous", "nous contacter",
        # Spanish
        "contacto", "contáctenos", "contáctanos",
        # Italian
        "contatto", "contattaci", "contatti"
    ]

    # Industry patterns - 40+ variations
    INDUSTRY_PATTERNS = [
        # English
        "industry", "industries", "sectors", "verticals", "markets",
        "industries we serve", "sectors we serve", "markets we serve",
        "industry focus", "sector focus", "market focus", "specialization",
        "our expertise", "areas of expertise", "domain expertise", "domains",
        "fields", "our fields", "business areas", "focus areas",
        # Common industries
        "technology", "healthcare", "finance", "fintech", "edtech", "healthtech",
        "e-commerce", "ecommerce", "retail", "manufacturing", "logistics",
        "real estate", "insurance", "banking", "automotive", "energy",
        "telecommunications", "media", "entertainment", "hospitality", "travel",
        "food", "agriculture", "pharma", "biotech", "legal", "hr", "marketing",
        "saas", "b2b", "b2c", "enterprise", "startup", "smb", "small business"
    ]

    # ==================== REGEX PATTERNS ====================

    # Email pattern
    EMAIL_PATTERN = r'\b[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Za-z]{2,}\b'

    # Phone patterns (international)
    PHONE_PATTERNS = [
        r'\+?\d{1,4}[-.\s]?\(?\d{1,4}\)?[-.\s]?\d{1,4}[-.\s]?\d{1,9}',
        r'\(\d{3}\)\s*\d{3}[-.\s]?\d{4}',
        r'\d{3}[-.\s]\d{3}[-.\s]\d{4}',
        r'\+\d{1,3}\s?\d{1,14}'
    ]

    # URL patterns
    WEBSITE_PATTERN = r'https?://(?:www\.)?[\w.-]+\.[a-z]{2,}(?:/[\w./-]*)?'
    DOMAIN_PATTERN = r'\b(?:www\.)?[\w-]+\.(?:com|org|net|io|co|ai|app|dev|tech|cloud|software|solutions|services)\b'

    # Social media patterns
    LINKEDIN_PATTERN = r'linkedin\.com/(?:company|in)/[\w-]+'
    TWITTER_PATTERN = r'(?:twitter|x)\.com/[\w-]+'
    FACEBOOK_PATTERN = r'facebook\.com/[\w.-]+'
    INSTAGRAM_PATTERN = r'instagram\.com/[\w.-]+'
    YOUTUBE_PATTERN = r'youtube\.com/(?:c/|channel/|@)?[\w-]+'

    # Pricing patterns - Multi-currency support
    PRICE_PATTERNS = [
        # USD
        r'\$\s?\d+(?:,\d{3})*(?:\.\d{2})?(?:\s*(?:/|per)\s*(?:mo(?:nth)?|yr|year|user|seat|license))?',
        # EUR
        r'€\s?\d+(?:,\d{3})*(?:\.\d{2})?(?:\s*(?:/|per)\s*(?:mo(?:nth)?|yr|year|user|seat|license))?',
        # GBP
        r'£\s?\d+(?:,\d{3})*(?:\.\d{2})?(?:\s*(?:/|per)\s*(?:mo(?:nth)?|yr|year|user|seat|license))?',
        # INR
        r'₹\s?\d+(?:,\d{3})*(?:\.\d{2})?(?:\s*(?:/|per)\s*(?:mo(?:nth)?|yr|year|user|seat|license))?',
        r'(?:INR|Rs\.?)\s?\d+(?:,\d{3})*(?:\.\d{2})?',
        # Generic with currency codes
        r'(?:USD|EUR|GBP|CAD|AUD)\s?\d+(?:,\d{3})*(?:\.\d{2})?',
        # Free tier
        r'\bfree\b(?:\s+(?:plan|tier|version|trial))?',
        # Custom/Contact
        r'\b(?:custom|contact\s+(?:us|sales)|get\s+quote|request\s+quote)\b'
    ]

    # Billing cycle patterns
    BILLING_PATTERNS = [
        (r'/\s*mo(?:nth)?|per\s+month|monthly|/m\b', 'monthly'),
        (r'/\s*yr|/\s*year|per\s+year|yearly|annually|/y\b', 'yearly'),
        (r'one[- ]time|lifetime|perpetual', 'one-time'),
        (r'/\s*user|per\s+user|per\s+seat', 'per-user'),
        (r'/\s*(?:hr|hour)|per\s+hour|hourly', 'hourly'),
        (r'/\s*(?:day)|per\s+day|daily', 'daily')
    ]

    # Plan tier names
    PLAN_NAMES = {
        'free': ['free', 'starter', 'basic', 'lite', 'trial', 'hobby'],
        'mid': ['pro', 'professional', 'standard', 'plus', 'growth', 'team', 'business'],
        'high': ['enterprise', 'premium', 'ultimate', 'unlimited', 'custom', 'corporate', 'advanced']
    }

    # Industry keywords for classification
    INDUSTRY_KEYWORDS = {
        'technology': ['software', 'saas', 'tech', 'digital', 'it', 'cloud', 'ai', 'ml', 'data'],
        'healthcare': ['health', 'medical', 'healthcare', 'pharma', 'biotech', 'clinical', 'patient'],
        'finance': ['finance', 'fintech', 'banking', 'insurance', 'investment', 'payment', 'trading'],
        'e-commerce': ['ecommerce', 'e-commerce', 'retail', 'shopping', 'marketplace', 'store'],
        'education': ['education', 'edtech', 'learning', 'training', 'course', 'school', 'university'],
        'marketing': ['marketing', 'advertising', 'seo', 'social media', 'content', 'branding'],
        'hr': ['hr', 'human resources', 'recruitment', 'hiring', 'talent', 'workforce'],
        'logistics': ['logistics', 'shipping', 'delivery', 'supply chain', 'warehouse', 'freight'],
        'real_estate': ['real estate', 'property', 'housing', 'rental', 'mortgage'],
        'legal': ['legal', 'law', 'attorney', 'compliance', 'contract']
    }

    # Common job titles for team extraction
    JOB_TITLES = [
        'ceo', 'cto', 'cfo', 'coo', 'cmo', 'cpo', 'cio', 'cso',
        'chief executive', 'chief technology', 'chief financial', 'chief operating',
        'chief marketing', 'chief product', 'chief information', 'chief strategy',
        'founder', 'co-founder', 'cofounder', 'president', 'vice president', 'vp',
        'director', 'head of', 'lead', 'manager', 'senior', 'principal',
        'partner', 'managing partner', 'general partner', 'advisor', 'board member'
    ]

    def __init__(self, fuzzy_threshold: int = 60):
        """
        Initialize parser with configurable threshold.

        Args:
            fuzzy_threshold: Minimum fuzzy match score (default: 60 for flexibility)
        """
        self.fuzzy_threshold = fuzzy_threshold

    def parse_file(self, file_content: bytes, filename: str) -> ParsedInfoDoc:
        """Parse info doc file (PDF, DOCX, or PPTX)."""
        logger.info(f"📄 [INFO DOC PARSER] Starting parse of: {filename}")

        try:
            filename_lower = filename.lower()

            if filename_lower.endswith('.pdf'):
                text = self._extract_text_from_pdf(file_content)
            elif filename_lower.endswith('.docx') or filename_lower.endswith('.doc'):
                text = self._extract_text_from_docx(file_content)
            elif filename_lower.endswith('.pptx') or filename_lower.endswith('.ppt'):
                text = self._extract_text_from_pptx(file_content)
            else:
                raise ValueError(f"Unsupported file format: {filename}")

            logger.info(f"📝 [INFO DOC PARSER] Extracted {len(text)} characters")
            return self._parse_text(text)

        except Exception as e:
            logger.error(f"❌ [INFO DOC PARSER] Failed to parse: {e}")
            return ParsedInfoDoc(warnings=[f"Failed to parse document: {str(e)}"])

    def _extract_text_from_pdf(self, file_content: bytes) -> str:
        """Extract text from PDF"""
        logger.info("📄 [PARSER] Extracting text from PDF...")
        text_parts = []

        with pdfplumber.open(io.BytesIO(file_content)) as pdf:
            for page_num, page in enumerate(pdf.pages, 1):
                page_text = page.extract_text()
                if page_text:
                    text_parts.append(page_text)
                    logger.debug(f"  Page {page_num}: {len(page_text)} chars")

                # Also extract tables
                tables = page.extract_tables()
                for table in tables:
                    for row in table:
                        if row:
                            row_text = " | ".join([str(cell) if cell else "" for cell in row])
                            if row_text.strip():
                                text_parts.append(row_text)

        return "\n".join(text_parts)

    def _extract_text_from_docx(self, file_content: bytes) -> str:
        """Extract text from DOCX"""
        logger.info("📄 [PARSER] Extracting text from DOCX...")
        doc = Document(io.BytesIO(file_content))
        text_parts = []

        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text_parts.append(paragraph.text)

        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join([cell.text.strip() for cell in row.cells if cell.text.strip()])
                if row_text:
                    text_parts.append(row_text)

        return "\n".join(text_parts)

    def _extract_text_from_pptx(self, file_content: bytes) -> str:
        """Extract text from PPTX"""
        logger.info("📄 [PARSER] Extracting text from PPTX...")

        try:
            from pptx import Presentation
            prs = Presentation(io.BytesIO(file_content))
            text_parts = []

            for slide_num, slide in enumerate(prs.slides, 1):
                slide_texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_texts.append(shape.text)
                    if shape.has_table:
                        for row in shape.table.rows:
                            row_text = " | ".join([cell.text.strip() for cell in row.cells if cell.text.strip()])
                            if row_text:
                                slide_texts.append(row_text)

                if slide_texts:
                    text_parts.append(f"--- Slide {slide_num} ---")
                    text_parts.extend(slide_texts)

            return "\n".join(text_parts)

        except ImportError:
            logger.warning("python-pptx not installed, trying basic extraction")
            return ""

    def _parse_text(self, text: str) -> ParsedInfoDoc:
        """Parse extracted text into structured data"""
        logger.info("🧠 [PARSER] Parsing text into sections...")

        info_doc = ParsedInfoDoc(raw_text=text)
        info_doc.word_count = len(text.split())
        lines = [line.strip() for line in text.split('\n') if line.strip()]

        # Extract contact info first (appears throughout document)
        info_doc.contact_info = self._extract_contact_info(text)

        # Extract company name (usually first prominent text)
        info_doc.company_name = self._extract_company_name(lines, text)

        # Extract tagline (often near company name)
        info_doc.tagline = self._extract_tagline(lines, text)

        # Detect industry
        info_doc.industry = self._detect_industry(text)

        # Detect sections
        sections = self._detect_sections(lines)
        info_doc.detected_sections = sections

        # Process each section
        for section in sections:
            content = section.content

            if section.section_type == "products":
                products = self._extract_products_services(content)
                info_doc.products_services.extend(products)

            elif section.section_type == "pricing":
                info_doc.pricing_raw = content
                tiers = self._extract_pricing_tiers(content)
                info_doc.pricing_tiers.extend(tiers)

            elif section.section_type == "benefits":
                benefits = self._extract_list_items(content)
                info_doc.key_benefits.extend(benefits)

            elif section.section_type == "usp":
                usps = self._extract_list_items(content)
                info_doc.unique_selling_points.extend(usps)

            elif section.section_type == "problem":
                info_doc.problem_solved = content[:1000]  # Limit length

            elif section.section_type == "team":
                members = self._extract_team_members(content)
                info_doc.team_members.extend(members)

        # Also scan full text for pricing if not found in sections
        if not info_doc.pricing_tiers:
            all_prices = self._find_all_prices(text)
            for price_str in all_prices[:10]:  # Limit to 10
                info_doc.pricing_tiers.append(PricingTier(
                    name="Detected Price",
                    price=price_str
                ))

        # Calculate confidence
        info_doc.confidence_score = self._calculate_confidence(info_doc)

        # Generate warnings
        if not info_doc.company_name:
            info_doc.warnings.append("Could not extract company name")
        if not info_doc.products_services:
            info_doc.warnings.append("No products/services detected")
        if not info_doc.contact_info.emails and not info_doc.contact_info.phones:
            info_doc.warnings.append("No contact information found")

        logger.info(
            f"✅ [PARSER] Parsing complete: "
            f"{len(info_doc.detected_sections)} sections, "
            f"{len(info_doc.products_services)} products/services, "
            f"confidence: {info_doc.confidence_score:.1f}%"
        )

        return info_doc

    def _extract_contact_info(self, text: str) -> ContactInfo:
        """Extract all contact information"""
        contact = ContactInfo()

        # Emails
        emails = re.findall(self.EMAIL_PATTERN, text, re.IGNORECASE)
        contact.emails = list(set(emails))[:10]  # Dedupe and limit

        # Phones
        for pattern in self.PHONE_PATTERNS:
            phones = re.findall(pattern, text)
            for phone in phones:
                if len(phone) >= 7 and phone not in contact.phones:
                    contact.phones.append(phone)
        contact.phones = contact.phones[:5]  # Limit

        # Websites
        urls = re.findall(self.WEBSITE_PATTERN, text, re.IGNORECASE)
        domains = re.findall(self.DOMAIN_PATTERN, text, re.IGNORECASE)

        all_sites = list(set(urls + ['https://' + d if not d.startswith('http') else d for d in domains]))
        # Filter out social media
        contact.websites = [
            url for url in all_sites
            if not any(social in url.lower() for social in ['linkedin', 'twitter', 'facebook', 'instagram', 'youtube'])
        ][:5]

        # Social media
        linkedin = re.search(self.LINKEDIN_PATTERN, text, re.IGNORECASE)
        if linkedin:
            contact.social_media['linkedin'] = 'https://' + linkedin.group(0)

        twitter = re.search(self.TWITTER_PATTERN, text, re.IGNORECASE)
        if twitter:
            contact.social_media['twitter'] = 'https://' + twitter.group(0)

        facebook = re.search(self.FACEBOOK_PATTERN, text, re.IGNORECASE)
        if facebook:
            contact.social_media['facebook'] = 'https://' + facebook.group(0)

        instagram = re.search(self.INSTAGRAM_PATTERN, text, re.IGNORECASE)
        if instagram:
            contact.social_media['instagram'] = 'https://' + instagram.group(0)

        youtube = re.search(self.YOUTUBE_PATTERN, text, re.IGNORECASE)
        if youtube:
            contact.social_media['youtube'] = 'https://' + youtube.group(0)

        return contact

    def _extract_company_name(self, lines: List[str], text: str) -> Optional[str]:
        """Extract company name from document"""
        # Strategy 1: Look for patterns like "Company Name Inc.", "ABC Corp", etc.
        company_patterns = [
            r'^([A-Z][A-Za-z0-9\s&]+(?:Inc\.?|Corp\.?|LLC|Ltd\.?|Limited|GmbH|AG|SA|SAS|BV|Pvt\.?\s*Ltd\.?))\s*$',
            r'^([A-Z][A-Za-z0-9\s&]{2,30})\s*$'  # Capitalized name at start
        ]

        for line in lines[:10]:  # Check first 10 lines
            line = line.strip()
            if 3 < len(line) < 50:
                for pattern in company_patterns:
                    match = re.match(pattern, line)
                    if match:
                        return match.group(1).strip()

        # Strategy 2: First capitalized line that's not a common header
        common_headers = {'about', 'services', 'products', 'contact', 'home', 'menu', 'pricing'}
        for line in lines[:15]:
            line_clean = line.strip()
            if 3 < len(line_clean) < 50 and line_clean.lower() not in common_headers:
                if line_clean[0].isupper() and not line_clean.isupper():
                    return line_clean

        return lines[0] if lines else None

    def _extract_tagline(self, lines: List[str], text: str) -> Optional[str]:
        """Extract company tagline/slogan"""
        # Taglines are usually short, impactful sentences
        tagline_patterns = [
            r'"([^"]{10,100})"',  # Quoted text
            r"'([^']{10,100})'",  # Single quoted
            r'([A-Z][^.!?]{10,80}[.!])',  # Short sentence
        ]

        # Look in first portion of text
        search_text = text[:2000]

        for pattern in tagline_patterns:
            matches = re.findall(pattern, search_text)
            for match in matches:
                # Filter out section headers and common phrases
                if len(match.split()) >= 3 and len(match.split()) <= 15:
                    return match.strip()

        # Look for line after company name
        for i, line in enumerate(lines[:10]):
            if len(line.split()) >= 3 and len(line.split()) <= 12:
                if not any(kw in line.lower() for kw in ['about', 'services', 'contact', 'copyright', 'privacy']):
                    return line

        return None

    def _detect_industry(self, text: str) -> Optional[str]:
        """Detect industry from document content"""
        text_lower = text.lower()

        industry_scores = {}
        for industry, keywords in self.INDUSTRY_KEYWORDS.items():
            score = sum(1 for kw in keywords if kw in text_lower)
            if score > 0:
                industry_scores[industry] = score

        if industry_scores:
            return max(industry_scores, key=industry_scores.get)

        return None

    def _detect_sections(self, lines: List[str]) -> List[InfoDocSection]:
        """Detect document sections using fuzzy matching"""
        sections = []

        section_patterns = {
            "company": self.COMPANY_PATTERNS,
            "products": self.PRODUCTS_PATTERNS,
            "pricing": self.PRICING_PATTERNS,
            "benefits": self.BENEFITS_PATTERNS,
            "usp": self.USP_PATTERNS,
            "problem": self.PROBLEM_PATTERNS,
            "team": self.TEAM_PATTERNS,
            "contact": self.CONTACT_PATTERNS,
            "industry": self.INDUSTRY_PATTERNS
        }

        i = 0
        while i < len(lines):
            line = lines[i]
            line_lower = line.lower().strip()

            # Skip very long lines (not headers)
            if len(line_lower) > 100:
                i += 1
                continue

            best_match = None
            best_score = 0
            best_section_type = None

            for section_type, patterns in section_patterns.items():
                for pattern in patterns:
                    # Exact match bonus
                    if line_lower == pattern:
                        score = 100
                    else:
                        score = fuzz.ratio(line_lower, pattern)

                    # Bonus for shorter lines (more likely headers)
                    if len(line_lower.split()) <= 4:
                        score += 10

                    if score > best_score:
                        best_score = score
                        best_match = pattern
                        best_section_type = section_type

            if best_score >= self.fuzzy_threshold:
                # Find section end
                section_start = i + 1
                section_end = len(lines)

                for j in range(i + 1, len(lines)):
                    next_line_lower = lines[j].lower().strip()

                    # Skip short lines
                    if len(next_line_lower) > 100:
                        continue

                    is_heading = False
                    for patterns in section_patterns.values():
                        for pattern in patterns:
                            if fuzz.ratio(next_line_lower, pattern) >= self.fuzzy_threshold:
                                section_end = j
                                is_heading = True
                                break
                        if is_heading:
                            break

                    if is_heading:
                        break

                content = "\n".join(lines[section_start:section_end]).strip()

                if content:
                    section = InfoDocSection(
                        section_type=best_section_type,
                        heading=line,
                        content=content,
                        confidence=min(100, best_score),
                        line_start=section_start,
                        line_end=section_end
                    )
                    sections.append(section)
                    logger.debug(f"  ✓ {best_section_type}: '{line}' ({best_score}%)")

                i = section_end
                continue

            i += 1

        return sections

    def _extract_products_services(self, content: str) -> List[ProductService]:
        """Extract products/services from section content"""
        products = []

        # Try to parse structured items
        items = self._extract_list_items(content)

        for item in items:
            # Try to extract name and description
            parts = re.split(r'[:\-–—]', item, maxsplit=1)
            name = parts[0].strip()
            description = parts[1].strip() if len(parts) > 1 else None

            # Look for price in item
            price = None
            for pattern in self.PRICE_PATTERNS:
                price_match = re.search(pattern, item, re.IGNORECASE)
                if price_match:
                    price = price_match.group(0)
                    break

            # Extract features (bullet points within item)
            features = []
            feature_patterns = re.findall(r'[•●○▪▸►]\s*([^•●○▪▸►\n]+)', item)
            features.extend([f.strip() for f in feature_patterns if f.strip()])

            # Determine category
            category = "service"
            if any(kw in item.lower() for kw in ['product', 'platform', 'software', 'tool', 'app']):
                category = "product"
            elif any(kw in item.lower() for kw in ['solution', 'system']):
                category = "solution"

            if name and len(name) > 2:
                products.append(ProductService(
                    name=name[:100],
                    description=description[:500] if description else None,
                    pricing=price,
                    features=features[:10],
                    category=category
                ))

        # Limit results
        return products[:20]

    def _extract_pricing_tiers(self, content: str) -> List[PricingTier]:
        """Extract pricing tiers from pricing section"""
        tiers = []

        # Split by plan names or price patterns
        lines = content.split('\n')
        current_tier = None
        current_features = []

        for line in lines:
            line_lower = line.lower().strip()

            # Check if this is a plan name
            is_plan_name = False
            for tier_type, names in self.PLAN_NAMES.items():
                if any(name in line_lower for name in names):
                    # Save previous tier
                    if current_tier:
                        current_tier.features = current_features[:15]
                        tiers.append(current_tier)

                    # Extract price from same line or nearby
                    price = None
                    for pattern in self.PRICE_PATTERNS:
                        price_match = re.search(pattern, line, re.IGNORECASE)
                        if price_match:
                            price = price_match.group(0)
                            break

                    # Detect billing cycle
                    billing_cycle = None
                    currency = None

                    for bc_pattern, bc_name in self.BILLING_PATTERNS:
                        if re.search(bc_pattern, line, re.IGNORECASE):
                            billing_cycle = bc_name
                            break

                    # Detect currency
                    if price:
                        if '$' in price:
                            currency = 'USD'
                        elif '€' in price:
                            currency = 'EUR'
                        elif '£' in price:
                            currency = 'GBP'
                        elif '₹' in price or 'INR' in price:
                            currency = 'INR'

                    current_tier = PricingTier(
                        name=line.strip()[:50],
                        price=price or "Contact for pricing",
                        billing_cycle=billing_cycle,
                        currency=currency,
                        is_popular=any(kw in line_lower for kw in ['popular', 'recommended', 'best value']),
                        is_enterprise=tier_type == 'high'
                    )
                    current_features = []
                    is_plan_name = True
                    break

            # If not a plan name, might be a feature
            if not is_plan_name and current_tier:
                # Clean up feature line
                feature = re.sub(r'^[•●○▪▸►✓✔\-\*]\s*', '', line).strip()
                if feature and len(feature) > 2 and len(feature) < 200:
                    current_features.append(feature)

        # Don't forget last tier
        if current_tier:
            current_tier.features = current_features[:15]
            tiers.append(current_tier)

        return tiers[:10]  # Limit to 10 tiers

    def _extract_team_members(self, content: str) -> List[TeamMember]:
        """Extract team members from team section"""
        members = []

        lines = content.split('\n')

        for i, line in enumerate(lines):
            line_lower = line.lower()

            # Look for job titles
            has_title = any(title in line_lower for title in self.JOB_TITLES)

            if has_title:
                # Try to extract name and title
                # Pattern: "Name - Title" or "Name, Title" or "Title: Name"
                patterns = [
                    r'^([A-Z][a-z]+(?:\s+[A-Z][a-z]+)+)\s*[-–—,]\s*(.+)$',
                    r'^(.+?):\s*([A-Z][a-z]+(?:\s+[A-Z][a-z]+)+)$'
                ]

                for pattern in patterns:
                    match = re.match(pattern, line.strip())
                    if match:
                        # Determine which group is name vs title
                        g1, g2 = match.groups()

                        if any(title in g1.lower() for title in self.JOB_TITLES):
                            name, role = g2, g1
                        else:
                            name, role = g1, g2

                        # Extract email if in nearby lines
                        email = None
                        linkedin = None

                        for check_line in lines[max(0, i-2):min(len(lines), i+3)]:
                            email_match = re.search(self.EMAIL_PATTERN, check_line)
                            if email_match:
                                email = email_match.group(0)

                            linkedin_match = re.search(self.LINKEDIN_PATTERN, check_line, re.IGNORECASE)
                            if linkedin_match:
                                linkedin = 'https://' + linkedin_match.group(0)

                        members.append(TeamMember(
                            name=name.strip()[:100],
                            role=role.strip()[:100],
                            title=role.strip()[:100],
                            email=email,
                            linkedin=linkedin
                        ))
                        break

        return members[:20]  # Limit to 20 members

    def _extract_list_items(self, content: str) -> List[str]:
        """Extract list items from content"""
        items = []

        # Pattern 1: Bullet points
        bullet_items = re.findall(r'[•●○▪▸►✓✔\-\*]\s*([^\n•●○▪▸►✓✔\-\*]+)', content)
        items.extend([item.strip() for item in bullet_items if item.strip()])

        # Pattern 2: Numbered items
        numbered_items = re.findall(r'\d+[\.\)]\s*([^\n\d]+)', content)
        items.extend([item.strip() for item in numbered_items if item.strip()])

        # Pattern 3: Lines that start with capital letter
        if not items:
            lines = content.split('\n')
            for line in lines:
                line = line.strip()
                if line and len(line) > 5 and line[0].isupper():
                    items.append(line)

        # Deduplicate while preserving order
        seen = set()
        unique_items = []
        for item in items:
            if item not in seen:
                seen.add(item)
                unique_items.append(item[:500])  # Limit length

        return unique_items[:30]  # Limit count

    def _find_all_prices(self, text: str) -> List[str]:
        """Find all prices in document"""
        prices = []

        for pattern in self.PRICE_PATTERNS:
            matches = re.findall(pattern, text, re.IGNORECASE)
            prices.extend(matches)

        # Deduplicate
        return list(set(prices))

    def _calculate_confidence(self, info_doc: ParsedInfoDoc) -> float:
        """Calculate parsing confidence score"""
        score = 0.0

        # Company info
        if info_doc.company_name:
            score += 15
        if info_doc.tagline:
            score += 5
        if info_doc.industry:
            score += 5

        # Products/Services
        if info_doc.products_services:
            score += 20
            if len(info_doc.products_services) >= 3:
                score += 5

        # Value proposition
        if info_doc.key_benefits:
            score += 10
        if info_doc.unique_selling_points:
            score += 10
        if info_doc.problem_solved:
            score += 5

        # Pricing
        if info_doc.pricing_tiers:
            score += 10
            if len(info_doc.pricing_tiers) >= 2:
                score += 5

        # Contact
        if info_doc.contact_info.emails:
            score += 5
        if info_doc.contact_info.phones:
            score += 3
        if info_doc.contact_info.websites:
            score += 2

        # Team
        if info_doc.team_members:
            score += 5

        return min(100.0, score)
